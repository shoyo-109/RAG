import os
import uuid
import logging
from typing import List
from .base_parser import BaseParserStrategy, ParseResult
try:
    from ..models import DocumentElement, ElementType, Provenance, CapabilityState, PipelineContext
except (ImportError, ValueError):
    from models import DocumentElement, ElementType, Provenance, CapabilityState, PipelineContext

logger = logging.getLogger("OfficeParserStrategy")


class OfficeParserStrategy(BaseParserStrategy):
    name = "OfficeParserStrategy"
    version = "1.0.0"

    def parse(self, file_path: str, context: PipelineContext) -> ParseResult:
        elements: List[DocumentElement] = []
        ext = os.path.splitext(file_path)[-1].lower()

        try:
            # Word documents (.docx)
            if ext == ".docx":
                import docx
                doc = docx.Document(file_path)
                for idx, p in enumerate(doc.paragraphs):
                    text = p.text.strip()
                    if not text:
                        continue
                    
                    is_heading = p.style.name.startswith("Heading") or idx == 0
                    el_type = ElementType.HEADING if is_heading else ElementType.PARAGRAPH

                    prov = Provenance(
                        source_file=file_path,
                        page=1,
                        parser_name="python-docx",
                        parser_version=getattr(docx, "__version__", "1.0")
                    )

                    elements.append(DocumentElement(
                        element_id=f"docx_el_{uuid.uuid4().hex[:8]}",
                        element_type=el_type,
                        text=text,
                        metadata={"style": p.style.name},
                        confidence=0.98,
                        provenance=prov
                    ))

                # Also extract tables in docx
                for t_idx, table in enumerate(doc.tables):
                    table_rows = []
                    for row in table.rows:
                        table_rows.append([cell.text.strip() for cell in row.cells])
                    table_str = "\n".join([" | ".join(r) for r in table_rows])
                    
                    if table_str.strip():
                        prov = Provenance(source_file=file_path, page=1, parser_name="python-docx")
                        elements.append(DocumentElement(
                            element_id=f"docx_tbl_{uuid.uuid4().hex[:8]}",
                            element_type=ElementType.TABLE,
                            text=table_str,
                            confidence=0.95,
                            provenance=prov
                        ))

            # PowerPoint (.pptx)
            elif ext == ".pptx":
                import pptx
                prs = pptx.Presentation(file_path)
                for slide_idx, slide in enumerate(prs.slides):
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            prov = Provenance(source_file=file_path, page=slide_idx + 1, parser_name="python-pptx")
                            elements.append(DocumentElement(
                                element_id=f"pptx_el_{uuid.uuid4().hex[:8]}",
                                element_type=ElementType.PARAGRAPH,
                                text=shape.text.strip(),
                                metadata={"slide": slide_idx + 1},
                                confidence=0.95,
                                provenance=prov
                            ))

            # Excel (.xlsx)
            elif ext in [".xlsx", ".xls"]:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, data_only=True)
                for sheetname in wb.sheetnames:
                    sheet = wb[sheetname]
                    sheet_data = []
                    for row in sheet.iter_rows(values_only=True):
                        row_vals = [str(val) if val is not None else "" for val in row]
                        if any(row_vals):
                            sheet_data.append(" | ".join(row_vals))

                    if sheet_data:
                        prov = Provenance(source_file=file_path, page=1, parser_name="openpyxl")
                        elements.append(DocumentElement(
                            element_id=f"xlsx_tbl_{uuid.uuid4().hex[:8]}",
                            element_type=ElementType.TABLE,
                            text=f"Sheet: {sheetname}\n" + "\n".join(sheet_data),
                            metadata={"sheet": sheetname},
                            confidence=0.95,
                            provenance=prov
                        ))

            if elements:
                return ParseResult(
                    elements=elements,
                    confidence=0.95,
                    status=CapabilityState.SUCCESS,
                    parser_name=self.name,
                    parser_version=self.version
                )

            return ParseResult(
                elements=[],
                confidence=0.0,
                status=CapabilityState.FAILED,
                error_message=f"No elements extracted from Office document {ext}",
                parser_name=self.name,
                parser_version=self.version
            )

        except Exception as e:
            logger.error(f"OfficeParserStrategy failed: {e}")
            return ParseResult(
                elements=[],
                confidence=0.0,
                status=CapabilityState.FAILED,
                error_message=str(e),
                parser_name=self.name,
                parser_version=self.version
            )
